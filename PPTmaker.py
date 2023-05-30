from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Define function to add slide
def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1] # Use 'Title and Content' layout (layout index 1)
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    # Add title and content to slide
    title_placeholder.text = title
    content_placeholder.text = content

# Slide 1
add_slide(prs,
          "The History of AI",
          "An overview of the evolution and milestones in Artificial Intelligence")

# Slide 2
add_slide(prs,
          "Early Beginnings",
          "The concept of AI has its roots in ancient history, with stories of artificial beings endowed with intelligence. However, the scientific journey began in the 1950s.")

# Slide 3
add_slide(prs,
          "The First Wave of AI",
          "From the 1950s through the 1970s, the first wave of AI was focused on rule-based systems and the development of the foundational AI concepts.")

# Slide 4
add_slide(prs,
          "AI Winter and Revival",
          "AI research faced a slowdown during the 1980s and 1990s due to various challenges. The revival started in the 2000s with advancements in machine learning and big data.")

# Slide 5
add_slide(prs,
          "Modern AI and the Future",
          "Today's AI leverages deep learning, natural language processing, and more. It continues to shape various industries and has a promising future.")

# Save presentation
prs.save('AI_History.pptx')
